import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Cm, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# Step 1: Import data from Excel
excel_file = 'dummy_data.xlsx'
raw_data = pd.read_excel(excel_file)

# Step 2: Transform data - convert the return columns to percentages
return_columns = [col for col in raw_data.columns if 'rtn' in col]  # Identify columns containing return data
raw_data[return_columns] = raw_data[return_columns] * 100  # Convert to percentage
raw_data[return_columns] = raw_data[return_columns].round(2)  # Round to two decimal places

# Function to estimate text width
def get_text_width(text, font_size=10):
    return len(text) * font_size * 0.08  # rough estimation, may need to adjust

# Step 3: Insert data into PowerPoint
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[5])  # Adding a slide with title and content layout

title = slide.shapes.title
title.text = "Fund Performance Data"

# Define the headers to match the desired output, including line breaks for "Rel"
headers = [
    "Fund Name", "Benchmark", "AM Ref", "Fund\n3M", "Rel\n3M", "Fund\n1-Yr", "Rel\n1-Yr",
    "Fund\n3-Yr (p.a)", "Rel\n3-Yr (p.a)", "Fund\n5-Yr (p.a)", "Rel\n5-Yr (p.a)"
]

# Map the headers to the corresponding columns in the DataFrame
column_mapping = {
    "Fund Name": "fund_name",
    "Benchmark": "benchmark",
    "AM Ref": "ref_code",
    "Fund\n3M": "3m_rtn_fund",
    "Rel\n3M": "3m_rtn_rel",
    "Fund\n1-Yr": "1yr_rtn_fund",
    "Rel\n1-Yr": "1yr_rtn_rel",
    "Fund\n3-Yr (p.a)": "3yr_rtn_fund",
    "Rel\n3-Yr (p.a)": "3yr_rtn_rel",
    "Fund\n5-Yr (p.a)": "5yr_rtn_fund",
    "Rel\n5-Yr (p.a)": "5yr_rtn_rel"
}

# Create table and populate it with data
table = slide.shapes.add_table(len(raw_data) + 1, len(headers), Inches(0.5), Inches(1.5), Inches(9), Inches(5)).table

# Set column names
for col_index, col_name in enumerate(headers):
    cell = table.cell(0, col_index)
    cell.text_frame.clear()  # Clear existing text frame
    # Add new paragraphs for header text with line breaks
    parts = col_name.split('\n')
    for i, part in enumerate(parts):
        if i == 0:
            p = cell.text_frame.paragraphs[0]
            p.text = part
        else:
            p = cell.text_frame.add_paragraph()
            p.text = part
        p.font.bold = True
        p.font.size = Pt(9)
        if col_index < 3:
            p.alignment = PP_ALIGN.LEFT
        else:
            p.alignment = PP_ALIGN.RIGHT
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE  # Center vertically
    table.rows[0].height = int(Cm(1.4))  # Convert to int

# Estimate width for the first three columns
col_widths = [0] * len(headers)
for row_index, row in raw_data.iterrows():
    for col_index in range(3):  # Only for the first three columns
        text_width = get_text_width(str(row[column_mapping[headers[col_index]]]), font_size=9)
        col_widths[col_index] = max(col_widths[col_index], text_width)

# Convert estimated widths to Cm
col_widths = [int(Cm(width)) for width in col_widths]

# Fixed width for remaining columns
fixed_width = int(Cm(1.3))

# Total width calculations
total_fixed_width = fixed_width * (len(headers) - 3)
total_auto_width = sum(col_widths[:3])
total_width = int(Cm(22.86))

# Adjust widths proportionally to fit into the total width
scaling_factor = (total_width - total_fixed_width) / total_auto_width
for i in range(3):
    col_widths[i] = int(col_widths[i] * scaling_factor)

# Set column widths in the table
for col_index in range(len(headers)):
    if col_index < 3:
        table.columns[col_index].width = col_widths[col_index]
    else:
        table.columns[col_index].width = fixed_width

# Fill table with data
for row_index, row in raw_data.iterrows():
    for col_index, header in enumerate(headers):
        column_name = column_mapping[header]
        if column_name:  # Only process if there's a corresponding column
            item = row[column_name]
            cell = table.cell(row_index + 1, col_index)
            if column_name in return_columns:  # Format return columns as percentage
                cell.text = f"{item:.1f}%"
            else:
                cell.text = str(item)
            cell.text_frame.paragraphs[0].font.size = Pt(8)
            if col_index < 3:  # Left align the first three columns
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
            else:  # Right align the remaining columns
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE  # Center vertically
    table.rows[row_index + 1].height = int(Cm(0.6))  # Convert to int

# Save the presentation
prs.save('test.pptx')